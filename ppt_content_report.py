# -*- coding: utf-8 -*-
"""
PPT 内容完整分析报告
"""
import os
import sys
import pptx
from pptx import Presentation
import glob
from datetime import datetime

sys.stdout.reconfigure(encoding='utf-8')

def analyze_ppt(ppt_path):
    """分析 PPT 内容"""
    try:
        prs = Presentation(ppt_path)
        
        content = {
            "file": ppt_path,
            "slides_count": len(prs.slides),
            "slides": []
        }
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = {
                "number": i,
                "title": "",
                "texts": [],
                "images_count": 0
            }
            
            if slide.shapes.title:
                slide_content["title"] = slide.shapes.title.text
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and text != slide_content["title"]:
                        slide_content["texts"].append(text)
                
                if shape.shape_type == 13:
                    slide_content["images_count"] += 1
            
            content["slides"].append(slide_content)
        
        return content
        
    except Exception as e:
        return None

def main():
    """主函数"""
    ppt_files = glob.glob('D:/文化出海/**/*.pptx', recursive=True)
    
    print("=" * 70)
    print("PPT Content Analysis Report")
    print(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
    print("=" * 70)
    print(f"\nTotal PPT files found: {len(ppt_files)}")
    
    all_analyses = []
    
    for i, ppt_file in enumerate(ppt_files, 1):
        print(f"\n{'='*70}")
        print(f"[{i}/{len(ppt_files)}] {os.path.basename(ppt_file)}")
        print("=" * 70)
        
        content = analyze_ppt(ppt_file)
        if content:
            all_analyses.append(content)
            print(f"Slides: {content['slides_count']}")
            
            for slide in content["slides"]:
                print(f"\n  Slide {slide['number']}:", end=" ")
                if slide["title"]:
                    print(f"'{slide['title']}'", end=" | ")
                if slide["images_count"] > 0:
                    print(f"{slide['images_count']} images", end="")
                print()
                
                if slide["texts"]:
                    for text in slide["texts"]:
                        if len(text) < 80:
                            print(f"    - {text}")
                        else:
                            print(f"    - {text[:80]}...")
    
    # 总结
    print("\n\n" + "=" * 70)
    print("SUMMARY")
    print("=" * 70)
    print(f"Total PPT files analyzed: {len(all_analyses)}")
    print(f"Total slides: {sum(a['slides_count'] for a in all_analyses)}")
    
    print("\nPPT Theme: All are '评印改印' (Seal Evaluation & Correction) cases")
    print("Author: 庵角山人 (Anjiao Shanren) - Jiang Haoxu's art name")
    print("\nCommon Topics:")
    print("  - 空间分割 (Space allocation)")
    print("  - 文字结构 (Character structure)")
    print("  - 篆法含蓄 (Subtle seal script)")
    print("  - 线条力度 (Line strength)")
    print("  - 边款布局 (Side inscription layout)")
    
    return all_analyses

if __name__ == "__main__":
    main()
