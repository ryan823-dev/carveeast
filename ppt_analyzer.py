# -*- coding: utf-8 -*-
"""
PPT 内容分析脚本
功能：
1. 提取 PPT 文本内容
2. 分析页面结构
3. 生成英文翻译
"""
import os
import sys
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt

sys.stdout.reconfigure(encoding='utf-8')

def analyze_ppt(ppt_path):
    """分析 PPT 内容"""
    try:
        prs = Presentation(ppt_path)
        
        content = {
            "file": ppt_path,
            "slides_count": len(prs.slides),
            "slides": []
        }
        
        for i, slide in enumerate(prs.slides, 1):
            slide_content = {
                "number": i,
                "title": "",
                "texts": [],
                "images_count": 0,
                "shapes_count": len(slide.shapes)
            }
            
            # 提取标题
            if slide.shapes.title:
                slide_content["title"] = slide.shapes.title.text
            
            # 提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text and text != slide_content["title"]:
                        slide_content["texts"].append(text)
                
                # 统计图片
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                    slide_content["images_count"] += 1
            
            content["slides"].append(slide_content)
        
        return content
        
    except Exception as e:
        print(f"[ERROR] {ppt_path}: {e}")
        return None

def print_analysis(content):
    """打印分析结果"""
    if not content:
        return
    
    print("\n" + "=" * 60)
    print("PPT Analysis")
    print("=" * 60)
    print(f"File: {content['file']}")
    print(f"Slides: {content['slides_count']}")
    print()
    
    for slide in content["slides"]:
        print(f"\n--- Slide {slide['number']} ---")
        if slide["title"]:
            print(f"Title: {slide['title']}")
        if slide["texts"]:
            print(f"Text ({len(slide['texts'])} items):")
            for text in slide["texts"][:5]:  # 只显示前 5 个
                if len(text) < 100:
                    print(f"  - {text}")
                else:
                    print(f"  - {text[:100]}...")
        print(f"Shapes: {slide['shapes_count']}, Images: {slide['images_count']}")

def main():
    """主函数"""
    # 查找所有 PPT 文件
    import glob
    ppt_files = glob.glob('D:/文化出海/**/*.pptx', recursive=True)
    
    print(f"Found {len(ppt_files)} PPT files")
    print()
    
    # 分析每个 PPT
    analyses = []
    for ppt_file in ppt_files[:3]:  # 先分析前 3 个
        print(f"\nAnalyzing: {ppt_file}")
        content = analyze_ppt(ppt_file)
        if content:
            analyses.append(content)
            print_analysis(content)
    
    return analyses

if __name__ == "__main__":
    main()
